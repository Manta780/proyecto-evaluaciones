import fitz  # PDF
from docx import Document  # Word
from pptx import Presentation  # PowerPoint
from io import BytesIO

from services.ocr_service import ocr_imagen, ocr_pdf
from services.ia_cleaner_service import limpieza_completa, limpieza_rapida



def extract_text(file_bytes: bytes, filename: str, modo_limpieza: str = "completa"):
    """
    Extrae texto de un archivo y opcionalmente lo limpia.

    Args:
        file_bytes: contenido del archivo
        filename: nombre del archivo (determina el tipo)
        modo_limpieza: "completa" (todo), "rapida" (solo caracteres), "ninguna" (sin limpieza)
    """
    filename = filename.lower()

    if filename.endswith(".pdf"):
        texto = extract_pdf(file_bytes)

        if not texto.strip():
            print("ADVERTENCIA: PDF sin texto, aplicando OCR...")
            texto = ocr_pdf(file_bytes)

        if modo_limpieza == "completa" and texto.strip():
            texto = limpieza_completa(texto)
        elif modo_limpieza == "rapida" and texto.strip():
            texto = limpieza_rapida(texto)

        return texto

    elif filename.endswith(".docx"):
        texto = extract_docx(file_bytes)
        if modo_limpieza == "completa" and texto.strip():
            texto = limpieza_completa(texto)
        elif modo_limpieza == "rapida" and texto.strip():
            texto = limpieza_rapida(texto)
        return texto

    elif filename.endswith(".pptx"):
        texto = extract_pptx(file_bytes)
        if modo_limpieza == "completa" and texto.strip():
            texto = limpieza_completa(texto)
        elif modo_limpieza == "rapida" and texto.strip():
            texto = limpieza_rapida(texto)
        return texto

    elif filename.endswith((".png", ".jpg", ".jpeg")):
        print("Aplicando OCR a imagen...")
        texto = ocr_imagen(file_bytes)
        if modo_limpieza == "completa" and texto.strip():
            texto = limpieza_completa(texto)
        elif modo_limpieza == "rapida" and texto.strip():
            texto = limpieza_rapida(texto)
        return texto

    else:
        raise ValueError("Formato no soportado")

# 📄 PDF
def extract_pdf(file_bytes: bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()
    return text


# 📝 WORD
def extract_docx(file_bytes: bytes):
    from io import BytesIO
    doc = Document(BytesIO(file_bytes))
    return "\n".join([p.text for p in doc.paragraphs])


# 📊 POWERPOINT
def extract_pptx(file_bytes: bytes):
    from io import BytesIO
    prs = Presentation(BytesIO(file_bytes))
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)